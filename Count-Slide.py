import collections.abc
from pptx import Presentation
import os


def find_files(root_dir, extension):
    result = []
    for root, dirs, files in os.walk(root_dir):
        # Add "extension" extension from "files" to "result".
        result.extend([os.path.join(root, f) for f in files if f.endswith(extension)])

    return result


if __name__ == '__main__':
    s = input("Dir: ")
    root_dir = s.strip('\"')
    files = find_files(root_dir, ".pptx")
    
    for file in files:
        print(file)
        prs = Presentation(file)

        # Number of slide
        print(f"Slide: {len(prs.slides)}")

        # Number of text boxes
        textbox_count = 0
        for slide in prs.slides:
            textbox_count += len(slide.shapes)

        print("Textbox:", textbox_count)

        # Number of slide masters
        print(f"SlideMaster: {len(prs.slide_masters)}")

        # Number of layout masters
        layout_master_count = 0
        for i, slide_master in enumerate(prs.slide_masters):
            layout_master_count += len(slide_master.slide_layouts)

        print("LayoutMaster:", layout_master_count)
        print("\n")